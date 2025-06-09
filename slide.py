import os
import requests
from dotenv import load_dotenv
from serpapi import GoogleSearch
from pptx import Presentation
from pptx.util import Inches, Pt

# Load API keys
load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
SERPAPI_API_KEY = os.getenv("SERPAPI_API_KEY")

# Web Search
def search_web(topic):
    params = {
        "q": topic,
        "api_key": SERPAPI_API_KEY,
        "engine": "google",
        "num": 5,
    }
    search = GoogleSearch(params)
    results = search.get_dict()
    snippets = [res["snippet"] for res in results.get("organic_results", []) if "snippet" in res]
    return "\n".join(snippets)

# Generate Slides from LLM (OpenRouter API)
def query_llm(prompt):
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "HTTP-Referer": "https://yourdomain.com",  # can be anything for personal use
    }
    body = {
        "model": "mistralai/mistral-7b-instruct",  # you can change this to gpt-3.5 or llama3
        "messages": [
            {"role": "system", "content": "You are a helpful assistant that creates PowerPoint slides."},
            {"role": "user", "content": prompt}
        ]
    }

    try:
        response = requests.post(url, headers=headers, json=body)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        return f"❌ LLM Error: {e}\nResponse: {response.text if response else 'No response'}"

# Create PowerPoint
def create_ppt_from_text(text, filename="slides.pptx"):
    prs = Presentation()
    slides = text.split("\n\n")

    for slide in slides:
        parts = slide.strip().split("\n", 1)
        if len(parts) < 2:
            continue
        title, content = parts
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title.strip()
        slide.placeholders[1].text = content.strip()

    prs.save(filename)
    print(f"✅ Presentation saved to {filename}")

#Generate Slide Content Prompt
def generate_slide_content(topic, web_text):
    return query_llm(
        f"Create 5 PowerPoint slides on the topic '{topic}'. Use the following info:\n{web_text}\n"
        "Format as:\n\nSlide Title\nBullet point 1\nBullet point 2\n...\n\n"
    )

#Main
def main():
    topic = input("Enter a topic: ")
    print("🔍 Searching the web...")
    web_text = search_web(topic)
    print("🧠 Generating slides with LLM...")
    slides_text = generate_slide_content(topic, web_text)
    print("📝 LLM Response:\n", slides_text)
    filename = f"{topic.replace(' ', '_')}_slides.pptx"
    print("📊 Creating presentation...")
    create_ppt_from_text(slides_text, filename)

if __name__ == "__main__":
    main()
